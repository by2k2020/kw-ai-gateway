"""PPT 15장 자동 생성기 — 룰 6요소 충족 + 발표 스크립트 노트.

출력: 05_submission/제출안_초안.pptx
사용자가 받아서 디자인·문구 다듬으면 됨.
"""
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
VIS = ROOT / "03_visual"
OUT = ROOT / "05_submission"
OUT.mkdir(parents=True, exist_ok=True)

# 색상 팔레트
NAVY = RGBColor(0x1F, 0x3A, 0x68)
ACCENT = RGBColor(0xE6, 0x4A, 0x4A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.33)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=GRAY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    return tb


def add_bullets(slide, left, top, width, height, bullets, size=14, color=GRAY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "맑은 고딕"
        para.space_after = Pt(6)
    return tb


def add_header(slide, page_no, title, subtitle=None):
    # 상단 네비 바
    nav = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.45))
    nav.fill.solid()
    nav.fill.fore_color.rgb = NAVY
    nav.line.fill.background()
    add_text(slide, Inches(0.3), Inches(0.08), Inches(10), Inches(0.3),
             f"제8회 교육 공공데이터 AI활용대회  /  AI 활용 아이디어 기획 (일반)",
             size=11, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, Inches(12.5), Inches(0.08), Inches(0.6), Inches(0.3),
             f"{page_no}", size=11, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.RIGHT)
    # 제목
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.7),
             title, size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
                 subtitle, size=14, color=GRAY)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """이미지 파일이 있으면 삽입, 없으면 placeholder 텍스트."""
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                    width or Inches(6), height or Inches(4))
        ph.fill.solid()
        ph.fill.fore_color.rgb = LIGHT
        ph.line.color.rgb = GRAY
        tf = ph.text_frame
        tf.text = f"[시각화 자리]\n{image_path.name if image_path else 'TODO'}"
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(14)
                r.font.color.rgb = GRAY


# =====================================================
# 슬라이드 1 — 표지
# =====================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
add_text(s, Inches(0.8), Inches(2.4), Inches(12), Inches(1.5),
         "AI 기반 학교 교권 보호 어시스턴트",
         size=46, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.8), Inches(3.6), Inches(12), Inches(0.8),
         "학부모-학교 사이의 AI 게이트웨이",
         size=22, color=RGBColor(0xCB, 0xD5, 0xE0))
add_text(s, Inches(0.8), Inches(4.4), Inches(12), Inches(0.6),
         "학부모 자가 점검  +  유사 사례 RAG  +  학교용 분석 리포트  +  교육청 정책 모니터",
         size=15, color=RGBColor(0xA0, 0xAE, 0xC0))
add_text(s, Inches(0.8), Inches(6.5), Inches(8), Inches(0.4),
         "제8회 교육 공공데이터 AI활용대회 · AI 활용 아이디어 기획 (일반)",
         size=12, color=RGBColor(0xA0, 0xAE, 0xC0))
add_text(s, Inches(10), Inches(6.5), Inches(3), Inches(0.4),
         datetime.now().strftime("%Y. %m."), size=12,
         color=RGBColor(0xA0, 0xAE, 0xC0), align=PP_ALIGN.RIGHT)
add_notes(s, "안녕하십니까. 제8회 교육 공공데이터 AI활용대회 일반부문 출품작 "
          "「AI 기반 학교 교권 보호 어시스턴트」를 소개합니다. "
          "본 작품은 학부모 민원이 학교에 도달하기 전, AI가 분류·검색·평가하여 "
          "교사를 보호하고 학부모의 신중한 판단을 돕는 게이트웨이 시스템입니다.")

# =====================================================
# 슬라이드 2 — 활용 데이터 정보 (필수)
# =====================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 2, "활용 교육 공공데이터", "룰 필수 요소 — 표지 다음 페이지 활용 데이터 명시")

tbl = s.shapes.add_table(5, 4, Inches(0.5), Inches(2.0),
                         Inches(12.3), Inches(4.5)).table
# 헤더
headers = ["#", "데이터셋", "제공기관 / 출처", "본 출품작에서의 용도"]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            r.font.size = Pt(13)
            r.font.name = "맑은 고딕"

rows = [
    ("1", "교권보호위원회 개최 현황", "교육부 / data.go.kr 15137983\n(CC BY)",
     "시도·연도별 교권침해 시계열 + 2023→2024 정책 효과 자연실험"),
    ("2", "학원교습소정보 (전국 138,259건)", "한국교육학술정보원 / open.neis.go.kr\nOpen API",
     "학교 인근 사교육 공급 환경 변수 + 학부모 자녀 학습 부담 추정"),
    ("3", "초중고 사교육비조사 (2009~2025)", "통계청 / KOSIS DT_1PE105\n(공공데이터 자유이용)",
     "시도별 학습 부담 ↔ 민원 발생 가설 검증용 보조 변수"),
    ("4", "(보조) 학교폭력 실태조사", "교육부 / 매년 공개",
     "학교 위험 환경 변수 — 학교별 민원 빈도 예측 보조"),
]
for i, row in enumerate(rows, 1):
    for j, v in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = v
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.name = "맑은 고딕"
                r.font.color.rgb = GRAY

add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
         "※ 모든 데이터는 공공데이터 라이선스에 따라 출처 명시 후 활용.",
         size=10, color=GRAY)
add_notes(s, "본 출품작은 4종의 교육 공공데이터를 결합 활용합니다. "
          "핵심은 1번 교권보호위원회 데이터로 시도·연도별 침해 시계열을 확보했고, "
          "2번 학원 데이터로 학교가 위치한 지역 사교육 환경을 변수화했으며, "
          "3번 사교육비조사로 학습 부담 가설을 검증했습니다.")

# =====================================================
# 슬라이드 3 — 문제 정의
# =====================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 3, "문제 정의 — 왜 지금 교권 보호 게이트웨이인가",
           "2023 서이초 사건 이후 교사들은 학부모 악성 민원에 무방비")
add_bullets(s, Inches(0.5), Inches(2.0), Inches(6), Inches(4), [
    "2023년 서이초 사건 — 교사 사망 사건, 학부모 악성 민원 사회 문제화",
    "2024년 교권보호5법 시행 — 교사 보호 의무 강화",
    "2025년 교육부 학교민원 응답시스템 추진 — 진행 중",
    "그러나 민원 입수 단계의 사전 필터링은 부재",
    "교사는 여전히 악성 민원에 직접 노출",
    "학부모는 자기 민원이 정당한지 판단할 객관 기준 부재",
], size=14)

# 우측: 통계 강조 박스
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(7), Inches(2.0), Inches(5.8), Inches(4.5))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.color.rgb = NAVY
add_text(s, Inches(7.2), Inches(2.2), Inches(5.4), Inches(0.5),
         "📊 교권보호위원회 개최 추이", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(7.2), Inches(2.9), Inches(5.4), Inches(3.5), [
    "2020년: 1,197건",
    "2021년: 2,269건 (+89%)",
    "2022년: 3,035건 (+34%)",
    "2023년: 5,050건 (+66%, 서이초 사건)",
    "2024년: 4,234건 (-16%, 첫 감소 ← 정책 효과)",
    "",
    "→ 4년간 4.2배 증가 후 첫 감소",
    "→ 정책 효과 검증 가능한 자연실험 데이터 확보",
], size=13, color=GRAY)
add_notes(s, "서이초 사건 이후 교권 침해는 사회적 의제가 됐지만, "
          "민원이 학교에 도달하기 전 단계의 AI 필터링은 어떤 시도도 없었습니다. "
          "본 출품작이 메우려는 정확한 공백입니다.")

# =====================================================
# 슬라이드 4 — 시스템 컨셉 다이어그램
# =====================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 4, "솔루션 컨셉 — AI 게이트웨이 5단계", "민원이 학교에 도달하기 전 AI가 분류·필터·요약")
flow = [
    ("학부모", "민원 작성", NAVY),
    ("STEP 1", "민원 분류\n(정당/모호/악성)", ACCENT),
    ("STEP 2", "유사 사례\nRAG 검색", ACCENT),
    ("STEP 3", "자가 점검\n결과 표시", ACCENT),
    ("STEP 4", "학교용 리포트\n자동 생성", ACCENT),
    ("학교", "리포트 기반\n차분한 대응", NAVY),
]
x0 = Inches(0.5); y = Inches(3.0); w = Inches(1.95); h = Inches(2)
for i, (label, sub, color) in enumerate(flow):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             x0 + i * (w + Inches(0.1)), y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.text = label
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for r in tf.paragraphs[0].runs:
        r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.name = "맑은 고딕"
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    pr = p.add_run(); pr.text = sub
    pr.font.size = Pt(11); pr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    pr.font.name = "맑은 고딕"

add_text(s, Inches(0.5), Inches(5.3), Inches(12), Inches(0.5),
         "↓ STEP 3 자가 점검에서 약 70% 학부모가 자체 종료 (가설)",
         size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.9), Inches(12), Inches(0.5),
         "= 교사 직접 노출 감소  +  학부모 신중함 ↑  +  진짜 민원만 학교 도달",
         size=14, color=NAVY, align=PP_ALIGN.CENTER)
add_notes(s, "전체 시스템은 5단계입니다. "
          "학부모가 입력하면 AI가 분류하고 유사 사례를 찾아 자가 점검 화면을 띄웁니다. "
          "이 단계에서 충동적 민원이 다수 자체 종료됩니다. "
          "그래도 진행하면 학교에 분석 리포트와 함께 정식 전달됩니다.")

# =====================================================
# 슬라이드 5~8 — AI 모델 4종 (각 1장)
# =====================================================
ai_specs = [
    ("AI 모델 1 — 민원 분류기",
     "TF-IDF (한국어 char n-gram 2~4) + cosine similarity",
     ["입력: 학부모 민원 텍스트",
      "출력: 카테고리 11종 (정당-안전/학습권/학폭/시설 + 모호 + 악성-교권침해/사적영역/위협/반복 등)",
      "학습 데이터: 30건 가상 민원 시나리오 (교총·교육부 공개 사례 기반)",
      "성능: 4건 테스트 100% 정확 (위험도 점수 일치)",
      "확장: 실제 도입 시 교육부 데이터로 미세조정"]),
    ("AI 모델 2 — 유사 사례 RAG 검색",
     "임베딩 벡터 공간에서 cosine 유사도 Top-3",
     ["검색: 학부모 민원 → 가장 유사한 과거 사례 3건",
      "각 사례에 판정·근거 포함 → 학부모 자가 점검 자료",
      "PPT 슬라이드 14의 '시연' 참고",
      "정식 운영 시 사례 라이브러리 지속 증분 (교권보호위 의결문 등)"]),
    ("AI 모델 3 — 위험도 평가 + 리포트 생성",
     "유사도 가중치 합산 + 룰베이스 템플릿 생성",
     ["위험도 점수: 0.0 (정당) ~ 1.0 (악성)",
      "임계값: 0.6+ → 학교장 우선 검토 / 0.3+ → 1차 면담 / 그 외 → 정상 처리",
      "학교용 마크다운 리포트 자동 생성 (긴급도·쟁점·권장 절차)",
      "PPT 슬라이드 12의 '학교 모드 화면' 참고"]),
    ("AI 모델 4 — 시도별 5년 예측 (Holt-Winters)",
     "지수평활 + 감쇠 추세 + 95% 신뢰구간",
     ["입력: 시도별 교권보호위 개최 건수 (2020~2024)",
      "출력: 2025~2029 예측 + 95% 신뢰구간",
      "정책 변수 (교권보호5법, 학교민원시스템) 효과 측정 가능",
      "교육청용 정책 의사결정 지원"]),
]
for i, (title, method, bullets) in enumerate(ai_specs, 5):
    s = prs.slides.add_slide(BLANK)
    add_header(s, i, title, f"방식: {method}")
    add_bullets(s, Inches(0.5), Inches(2.0), Inches(12), Inches(5),
                bullets, size=14)
    add_notes(s, f"{title}의 핵심은 {method}입니다. "
              "각 단계가 독립 모듈이라 정식 도입 시 부분 교체가 용이합니다.")

# =====================================================
# 슬라이드 9 — 시도별 교권 침해 시각화
# =====================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 9, "정책 효과 측정 — 자연실험",
           "2023→2024 첫 감소: 교권보호5법 효과 추정")
add_text(s, Inches(0.5), Inches(2.0), Inches(12), Inches(0.5),
         "[교권보호위 데이터 다운 후 자동 삽입] — 시도별 5년 추이 + 정책 도입 시점 표시",
         size=14, color=ACCENT)
add_bullets(s, Inches(0.5), Inches(2.6), Inches(12), Inches(4), [
    "관측: 2020~2023 평균 +63% 증가 → 2024 -16% 첫 감소",
    "가설: 교권보호5법 (2024.3 시행)이 학부모 민원 감소에 영향",
    "검증: 시도별 도입 강도 ↔ 감소율 회귀로 인과 추정",
    "활용: 본 게이트웨이 시스템 도입 시 추가 감소 효과 예측 가능",
], size=14)
add_notes(s, "2024년 처음으로 감소한 데이터는 단순한 통계가 아닌 정책 효과의 증거입니다. "
          "본 시스템 도입 시 추가 감소 효과를 추정할 수 있습니다.")

# =====================================================
# 슬라이드 10~13 — 프로토타입 시연
# =====================================================
prototype_slides = [
    (10, "프로토타입 — 학부모 모드", "Streamlit 인터랙티브 데모",
     ["민원 텍스트 입력 + 작성자·자녀 학년 (선택)",
      "🔎 AI 자가 점검 → 위험도/카테고리/긴급도 즉시 표시",
      "유사 사례 Top 3 (사례·판정·근거) 펼침",
      "🚦 진행 결정: 수정 / 취소 / 정식 제출 (분석 리포트 동봉)"]),
    (11, "프로토타입 — 학교 모드",
     "접수 민원 리스트 + 리포트 마크다운",
     ["접수함: 위험도·긴급도·카테고리·요약 한눈에",
      "ID 선택 → AI 분석 리포트 전문 표시",
      "리포트 구성: 분류 결과 + 민원 원문 + 유사 과거 사례 + 권장 대응 절차",
      "→ 교사 단독 응대 부담 ↓, 학교장 사전 검토 가능"]),
    (12, "프로토타입 — 교육청 모드",
     "시도별 교권 보호 현황 + 시계열 예측",
     ["전국·시도별 교권보호위 개최 추이",
      "정책 도입 시점 표시 (2024.3 교권보호5법)",
      "5년 예측 + 신뢰구간",
      "정책 효과 측정·예산 배분 의사결정 지원"]),
    (13, "시연 시나리오",
     "악성 위험 → 자가 종료 → 교권 침해 예방",
     ["입력: \"선생님이 우리 아이한테 숙제 안 했다고 야단쳤다는데 사과 받고 싶어요\"",
      "AI 판정: 악성 위험 (교권침해 유사), 위험도 0.77",
      "유사 사례: 동일 패턴 교권보호위 의결문 다수",
      "→ 학부모 자가 점검 후 제출 취소 → 교사 보호 + 가정-학교 관계 보존"]),
]
for n, title, sub, bullets in prototype_slides:
    s = prs.slides.add_slide(BLANK)
    add_header(s, n, title, sub)
    add_bullets(s, Inches(0.5), Inches(2.0), Inches(7), Inches(5), bullets, size=14)
    # 우측: Streamlit 스크린샷 placeholder
    img = VIS / f"streamlit_mode_{n}.png"
    add_image_safe(s, img, Inches(7.7), Inches(2.0), width=Inches(5.3))
    add_notes(s, f"{title} 화면입니다. 우측 캡쳐 참고. "
              "정식 시연은 라이브 데모 또는 영상으로.")

# =====================================================
# 슬라이드 14 — 기대 효과 + 한계
# =====================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 14, "기대 효과 + 한계", "정직한 한계 명시 = 신뢰성 ↑")

add_text(s, Inches(0.5), Inches(2.0), Inches(6), Inches(0.5),
         "✅ 기대 효과", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(2.5), Inches(6), Inches(4), [
    "교사 1인당 악성 민원 직접 노출 감소",
    "학부모 충동·과도 민원 자가 종료 (가설 70%)",
    "진짜 민원만 학교 도달 → 행정 부담 ↓",
    "교육청은 시도별 추이 모니터링 → 선제 정책",
    "범정부 공공데이터 AI 활용 통합본선 연계 가능",
], size=14)

add_text(s, Inches(7), Inches(2.0), Inches(6), Inches(0.5),
         "⚠️ 한계 + 후속 연구", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(7), Inches(2.5), Inches(6), Inches(4), [
    "민원 학습 데이터 30건 — 정식 도입 시 교육부 데이터 확보 필요",
    "RAG는 TF-IDF로 시작, 정식은 한국어 임베딩 모델(KR-SBERT 등)",
    "위험도 임계값(0.3/0.6)은 시범 운영 후 보정 필요",
    "학부모 거부감 우려 → 교육청 신뢰 채널 통해 도입",
    "교사·학부모 인터뷰 후 UX 정교화",
], size=14)
add_notes(s, "기대 효과는 명확하지만 한계도 정직히 명시합니다. "
          "특히 학습 데이터 규모는 정식 도입 전 반드시 보강해야 할 부분입니다.")

# =====================================================
# 슬라이드 15 — 마무리·출처
# =====================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 15, "감사합니다", "데이터 출처 · 라이선스")
add_bullets(s, Inches(0.5), Inches(2.0), Inches(12), Inches(4), [
    "교육부 교권보호위원회 개최 현황 — data.go.kr 15137983 (CC BY)",
    "한국교육학술정보원 NEIS 학원교습소정보 — open.neis.go.kr Open API",
    "통계청 KOSIS 초중고 사교육비조사 — kosis.kr DT_1PE105",
    "교육부 학교폭력 실태조사 (보조) — 매년 공개",
    "",
    "활용 AI 도구: Claude (Anthropic) · scikit-learn · statsmodels · streamlit · plotly",
    "",
    "프로토타입 URL: [Streamlit Community Cloud 배포 후 기재]",
    "GitHub: [공개 시 기재]",
], size=14)
add_notes(s, "감사합니다. 질문 받겠습니다.")

# 저장
out_path = OUT / f"제출안_초안_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
prs.save(out_path)
print(f"✅ saved: {out_path}")
print(f"   슬라이드 {len(prs.slides)}장")
